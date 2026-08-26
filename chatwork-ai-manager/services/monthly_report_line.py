"""業務月報の材料受付（TASK-20260826-002）。LINEを月報の入力源・トリガーにする。

これまでは鷲見がChatworkの全体ルームへ会議資料をアップロードしたことがトリガーだったが、
オーナー指示（2026-08-26）で**Chatwork起点を廃止**し、**オーナー本人がLINEで直接送った内容**
だけを月報の材料・トリガーにする（Chatworkの資料アップロードでは二度と月報を作らない。
services/scheduler.py からもChatwork検知の呼び出しを外した）。

UX（実装側で決定）:
  - オーナーがLINEで「月報開始」と送ると材料受付セッションを開く（既に開いていれば何もしない）。
  - **セッションが開いている間は、送った内容が通常のAI質問応答(qa.answer)を経由せず、
    そのまま材料として貯まる**（雑談・質問と混ざらないよう明示的に区切るため）。
    テキスト・画像（claude visionで内容を読む）・ファイル（PDF/Excel/Word/PowerPoint等）に対応。
  - 「月報終了」で締めて、貯まった材料からAIが月報を作る（Excel化・Dropbox保管・Chatwork
    アップ・社内メールは services/monthly_report.py 側で従来どおり行う。入力経路だけが変わる）。
  - セッションを開始していない通常のLINEメッセージは、従来どおり qa.answer に渡る。
  - **開いたまま放置される事故対策**: 一定時間（既定180分）操作が無い開いたセッションは
    scheduler.py が自動で締め切り、そこまでの材料で月報を作る（services/scheduler.py の
    run_monthly_report_line_check）。忘れて次の日に別件をLINEで聞いても、通常のQ&Aを
    ずっと乗っ取り続けることはない。
"""
import datetime
import os
import shutil
import tempfile

from db.connection import get_conn, query, query_one

START_WORDS = {"月報開始", "月報作成開始", "月報スタート"}
END_WORDS = {"月報終了", "月報作成終了", "月報完了", "月報終わり"}


def parse_command(text: str):
    """「月報開始」「月報終了」相当のコマンドかどうか。該当なければ None。"""
    t = (text or "").strip()
    if t in START_WORDS:
        return "start"
    if t in END_WORDS:
        return "end"
    return None


def current_session():
    row = query_one(
        "SELECT * FROM monthly_report_line_sessions WHERE status='open' ORDER BY id DESC LIMIT 1")
    return dict(row) if row else None


def start_session(line_user_id: str) -> dict:
    """既に開いていればそれを返す（冪等。二重に開始しない）。"""
    existing = current_session()
    if existing:
        return existing
    with get_conn() as conn:
        cur = conn.execute(
            "INSERT INTO monthly_report_line_sessions (line_user_id, status) VALUES (?, 'open')",
            (line_user_id,))
        session_id = cur.lastrowid
    return dict(query_one("SELECT * FROM monthly_report_line_sessions WHERE id=?", (session_id,)))


def close_session(session_id: int):
    with get_conn() as conn:
        conn.execute(
            "UPDATE monthly_report_line_sessions SET status='closed', closed_at=datetime('now') "
            "WHERE id=? AND status='open'", (session_id,))


def is_expired(session: dict) -> bool:
    from services import settings
    timeout_min = settings.get_int("monthly_report_line_session_timeout_min", 180)
    try:
        opened = datetime.datetime.strptime(session["opened_at"], "%Y-%m-%d %H:%M:%S")
    except (TypeError, ValueError):
        return False
    return (datetime.datetime.utcnow() - opened).total_seconds() > timeout_min * 60


def items(session_id: int) -> list:
    return [dict(r) for r in query(
        "SELECT * FROM monthly_report_line_items WHERE session_id=? ORDER BY id", (session_id,))]


def item_count(session_id: int) -> int:
    row = query_one(
        "SELECT COUNT(*) AS c FROM monthly_report_line_items WHERE session_id=?", (session_id,))
    return row["c"] if row else 0


def _add_item(session_id, kind, filename, text, ok, error):
    with get_conn() as conn:
        conn.execute(
            "INSERT INTO monthly_report_line_items "
            "(session_id, kind, filename, text, ok, error) VALUES (?, ?, ?, ?, ?, ?)",
            (session_id, kind, filename, text, 1 if ok else 0, error))


def _extract_pptx(path) -> str:
    """PowerPoint（会議資料の定番）は社内ナレッジ取込の抽出器に対応が無いので個別に読む。"""
    from pptx import Presentation
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        texts.append(" | ".join(cells))
        if texts:
            slides.append(f"[スライド{i}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


def _extract_line_file(message_id: str, filename: str):
    """LINEのファイルメッセージをダウンロードしテキスト化する。戻り値: (text_or_None, error_or_None)。

    PowerPointはここで専用抽出、それ以外は社内ナレッジと同じ抽出器(knowledge.extract)を使う
    （ダウンロード自体は services.attachments の実装を借用し、二重実装しない）。
    """
    from services import config, knowledge
    from services.attachments import MAX_BYTES, _download

    token = config.get("line_channel_access_token")
    if not token:
        return None, "LINEのアクセストークンが未設定"
    ext = os.path.splitext(filename or "")[1].lower()
    if ext not in (".ppt", ".pptx") and ext not in knowledge._EXTRACTORS:
        return None, f"未対応の形式（{ext or '拡張子なし'}）"
    tmp = tempfile.mkdtemp(prefix="mr-line-")
    try:
        path = os.path.join(tmp, os.path.basename(filename) or "line-file")
        _download(f"https://api-data.line.me/v2/bot/message/{message_id}/content", path,
                  headers={"Authorization": f"Bearer {token}"})
        if ext in (".ppt", ".pptx"):
            text = _extract_pptx(path)
        else:
            text = "\n\n".join(t for t, _ref in (knowledge.extract(path) or []) if t)
        return (text or None), (None if text else "本文を抽出できませんでした")
    except ValueError:
        return None, f"{MAX_BYTES // 1024 // 1024}MB を超えるので読みませんでした"
    except Exception as e:
        return None, f"{type(e).__name__}: {e}"
    finally:
        shutil.rmtree(tmp, ignore_errors=True)


def capture(session_id: int, msg: dict) -> str:
    """1件のLINEメッセージ（月報セッション中）を材料として保存し、オーナーへのack文を返す。"""
    mtype = msg.get("type")
    if mtype == "text":
        text = (msg.get("text") or "").strip()
        _add_item(session_id, "text", None, text, True, None)
    elif mtype == "image":
        from services.attachments import read_line_image
        text = read_line_image(msg.get("id"))
        _add_item(session_id, "image", "画像", text, True, None)
    elif mtype == "file":
        filename = msg.get("fileName") or "ファイル"
        text, error = _extract_line_file(msg.get("id"), filename)
        _add_item(session_id, "file", filename, text or "", text is not None, error)
    else:
        return "この種類のメッセージは月報の材料にできません（テキスト・画像・ファイルのみ対応）。"
    n = item_count(session_id)
    return (f"📥 月報の材料として受け付けました（現在 {n} 件）。\n"
            "続けて送るか、終わったら「月報終了」と送ってください。")
