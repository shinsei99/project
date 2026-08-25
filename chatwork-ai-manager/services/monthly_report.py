"""業務月報の自動作成（TASK-20260825-001）。

日報（社員ごと・毎日・実務の記録）とは役割も起点も別物:
  - 起点: 鷲見が**月1回**、全体Chatworkルームへ会議資料・会議内容をアップロードするメッセージ
    （時刻指定はない。トリガー駆動＝そのメッセージを検出したときに作る）。
  - 内容: 月次会議（社長挨拶→各担当の業務報告→今後の改善点・現状の問題点の3部構成）の記録。
    社長挨拶は月報に含めず、**各担当の業務報告**（鷲見のみ先頭固定・それ以降（塚本・松本・森・
    吉浦・大鹿）は会議によって順番が変わるため会議資料に記載された順（実際の発言順）に従う。
    欠席者は省略）と**今後の改善点・現状の問題点**の2セクションで出力する
    （TASK-20260825-003・順序の固定を解除＝TASK-20260825-005）。

トリガー判定:
  - 対象ルーム（monthly_report_room_id）で、対象者（monthly_report_account_id・既定は鷲見）本人が
    送った、**会議資料らしい拡張子のファイルが添付された**メッセージを「資料アップロード」とみなす。
    拡張子を絞るのは、鷲見さんが同じルームへ業者連絡先の写真等（.jpg/.png）を貼ることがあり、
    それを月報のトリガーにしないため（2026-08-25 実データで実際に IMG_*.JPG/.PNG が同居していた）。
  - 同じ日に複数の添付メッセージがあっても、月報は**日単位で1本**にする
    （その日最初のメッセージを正規のトリガーとし、他の添付・テキストは related_messages() でまとめて拾う）。
  - 1トリガー = 1本。monthly_reports.trigger_message_id の UNIQUE と scheduler._claim の二重で
    再生成しない（生成のやり直しは画面から明示的に行う）。
  - 出力（Excel）は日報と**同じ保存フォルダ**（daily_report_save_dir）に保管する。
"""
import datetime
import json
import os
import re
import tempfile

from db.connection import get_conn, query, query_one
from services import settings
from services.claude_client import run_json
from services.daily_report import id2name_map, readable

# Chatwork のファイル添付は本文に `[download:ID]ファイル名 (12.3 KB)` の形で入る。
# ここからファイル名まで取れるので、ダウンロード前に拡張子で絞り込める。
_RE_ATTACH = re.compile(r"\[download:(\d+)\](.+?)\s*\([\d.]+\s*[KMGkmg]?[Bb]\)")

# 会議資料として扱う拡張子（画像・音声・動画は含めない＝会議資料の想定外）
_MEETING_DOC_EXT = {".pdf", ".doc", ".docx", ".xls", ".xlsx", ".xlsm",
                    ".ppt", ".pptx", ".csv", ".txt", ".md", ".key", ".pages", ".numbers"}


# --- トリガー判定・対象メッセージの収集 -----------------------------------------
def target_room_id():
    rid = settings.get_setting("monthly_report_room_id", "") or \
        settings.get_setting("daily_report_room_id", "") or \
        settings.get_setting("manager_room_id", "")
    if rid:
        return int(rid)
    row = query_one("SELECT room_id FROM rooms WHERE monitored=1 AND type='group' "
                    "ORDER BY room_id LIMIT 1")
    return row["room_id"] if row else None


def trigger_account_id():
    v = settings.get_setting("monthly_report_account_id", "7426045")
    return int(v) if v else None


def _attached_docs(body: str) -> list:
    """本文から会議資料らしい添付だけを (file_id, filename) で拾う（画像等は除く）。"""
    out = []
    for fid, name in _RE_ATTACH.findall(body or ""):
        name = name.strip()
        if os.path.splitext(name)[1].lower() in _MEETING_DOC_EXT:
            out.append((int(fid), name))
    return out


def pending_triggers(limit: int = 5, room_id=None, account_id=None):
    """まだ月報にしていない「資料アップロード日」を、日ごとに1件（最初のメッセージ）で返す。"""
    rid = room_id if room_id is not None else target_room_id()
    aid = account_id if account_id is not None else trigger_account_id()
    if not rid or not aid:
        return []
    done = {r["trigger_message_id"] for r in query(
        "SELECT trigger_message_id FROM monthly_reports")}
    rows = query(
        "SELECT * FROM messages WHERE room_id=? AND account_id=? ORDER BY send_time, message_id",
        (rid, aid))
    by_day = {}
    for r in rows:
        r = dict(r)
        if not _attached_docs(r["body"]):
            continue
        day = datetime.date.fromtimestamp(r["send_time"]).isoformat()
        by_day.setdefault(day, r)   # 送信時刻順に並んでいるので、最初の1件がその日の正規トリガー
    out = [r for day, r in sorted(by_day.items()) if r["message_id"] not in done]
    return out[:limit]


def related_messages(trigger: dict):
    """トリガーと同じ日に、本人がそのルームへ送った発言（追加の資料・会議内容のテキスト）。"""
    d = datetime.date.fromtimestamp(trigger["send_time"])
    start = int(datetime.datetime.combine(d, datetime.time.min).timestamp())
    end = start + 86400
    rows = query(
        "SELECT * FROM messages WHERE room_id=? AND account_id=? AND send_time>=? AND send_time<? "
        "ORDER BY send_time, message_id",
        (trigger["room_id"], trigger["account_id"], start, end))
    return [dict(r) for r in rows]


# --- 添付資料のテキスト抽出 ------------------------------------------------------
def _extract_pptx(path) -> str:
    """PowerPoint（会議資料の定番）はナレッジ取込の対応外なので、ここだけ個別に読む。"""
    from pptx import Presentation
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        texts.append(" | ".join(cells))
        if texts:
            slides.append(f"[スライド{i}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


def extract_files_text(client, room_id, msgs) -> list:
    """会議資料らしい添付ファイルをダウンロードしテキストを抜き出す。

    戻り値: [(filename, text_or_None, error_or_None), ...]
    対応形式: PDF/Excel/Word/CSV/txt/md（社内ナレッジ取込と同じ）＋ PowerPoint。
    非対応形式・抽出失敗は本文を諦め、理由だけ残す（想像で埋めない）。
    """
    from services import knowledge
    out = []
    for m in msgs:
        for fid, name in _attached_docs(m["body"]):
            try:
                data, filename = client.download_file(room_id, fid)
                filename = filename or name
                if not data:
                    out.append((filename, None, "ダウンロードURLが取得できませんでした"))
                    continue
                ext = os.path.splitext(filename)[1].lower()
                with tempfile.TemporaryDirectory() as d:
                    path = os.path.join(d, filename)
                    with open(path, "wb") as f:
                        f.write(data)
                    if ext in (".ppt", ".pptx"):
                        text = _extract_pptx(path)
                    else:
                        parts = knowledge.extract(path)
                        text = "\n\n".join(t for t, _ref in parts if t)
                out.append((filename, text or None,
                           None if text else "対応形式でないか、本文を抽出できませんでした"))
            except Exception as e:
                out.append((name, None, f"{type(e).__name__}: {e}"))
    return out


# --- プロンプト --------------------------------------------------------------
_PROMPT = """あなたは不動産会社「大京商事」のAI業務マネージャーです。
{date}（{wd}曜日）に鷲見さんが全体Chatworkルームへアップロードした**会議資料・会議内容**から、
会社としての月報を書いてください。

# 会議の構成（月報はこれに合わせる）
実際の会議は次の3部構成です。
1. 社長挨拶 → 月報には**含めない**（読み取れても書かない）
2. 各担当の業務報告（**鷲見が1番目で固定**。それ以降の塚本・松本・森・吉浦・大鹿は
   会議によって順番が変わるため、会議資料・会議内容に記載された順（または実際の発言順）に従う。
   欠席者は報告なし）
3. 今後の改善点や現状の問題点

# 絶対に守ること
- **会議資料・会議内容に書かれていないことは書かない。** 想像で内容を作らない。
- 推測せざるを得ないことは文末に「（推測）」と付ける。
- 日報（社員ごとの日々の実務記録）とは役割が違います。月報は**会議で共有された内容**を
  まとめる位置づけです。
- 社長挨拶の内容（冒頭の挨拶・訓示など）は読み取れても本文に含めない。
- 添付資料のうち抽出に失敗しているものがあれば、内容を想像で埋めず、
  「〇〇（ファイル名）は内容を読み取れませんでした」とだけ書く。
- 会議内容のテキストが空で、添付資料の抽出もできなかった場合は、
  無理に本文を作らず「資料の内容を読み取れず、月報を作成できませんでした」とだけ書く。

# 鷲見さんの発言（会議内容のテキスト・{n_msgs}件）
{conversation}

# 添付資料から抽出したテキスト
{files}

# 出力
次のJSONだけを返してください（前後に説明文を付けない）。

{{
  "summary": "1行要約（40字以内）",
  "body_md": "月報本文（Markdown）"
}}

body_md は**次の2つの大見出し（##）だけ**を、この順で書いてください。見出しを増やさない。
各行は「・」ではなく `- ` で始める箇条書き。

## 各担当の業務報告

この下に、会議資料の中で実際に業務報告がある担当だけ、小見出し（###・担当者名のみ）を作り、
その人の報告内容を箇条書きでまとめる。**順番のルール:** 鷲見の報告があれば必ず1番目に置く。
それ以降（塚本・松本・森・吉浦・大鹿）は固定順にせず、会議資料に書かれた順（または実際に
発言・報告した順）のまま並べる。順番が資料から読み取れない担当者は、資料内での記載位置を
そのまま使う。
欠席等でその担当の報告が資料に無い場合は、その人の小見出しごと**丸ごと省略する**
（「特になし」と書かない・無理に埋めない）。6名とも報告が無ければ、この見出しの下に
「今回の資料からは各担当の業務報告を読み取れませんでした」と1行だけ書く。

## 今後の改善点や現状の問題点

会議で挙がった改善点・問題点を箇条書きでまとめる。該当が無ければ「特になし」と1行だけ書く。
"""


def build_prompt(trigger: dict, msgs: list, files_text: list) -> str:
    id2n = id2name_map()
    lines = []
    for m in msgs:
        t = datetime.datetime.fromtimestamp(m["send_time"]).strftime("%H:%M")
        text = readable(m["body"], id2n).replace("\n", "\n      ")
        lines.append(f"{t} {m['account_name']}: {text}")
    conversation = "\n".join(lines) if lines else "（この日のテキスト発言はありません）"

    blocks = []
    for filename, text, err in files_text:
        if text:
            blocks.append(f"### 添付資料: {filename}\n{text[:12000]}")
        else:
            blocks.append(f"### 添付資料: {filename}\n（{err}）")
    files_block = "\n\n".join(blocks) if blocks else "（添付資料はありません）"

    d = datetime.date.fromtimestamp(trigger["send_time"])
    wd = "月火水木金土日"[d.weekday()]
    return _PROMPT.format(date=d.isoformat(), wd=wd, n_msgs=len(msgs),
                          conversation=conversation, files=files_block)


# --- 生成・保存 ---------------------------------------------------------------
def model() -> str:
    return settings.get_setting("model_monthly_report", "sonnet")


def generate(trigger_message_id: str, generated_by: str = "manual", client=None) -> dict:
    """1本の月報を作って保存し、保存した行を返す。既存トリガーなら上書き（冪等）。"""
    trigger = query_one("SELECT * FROM messages WHERE message_id=?", (trigger_message_id,))
    if not trigger:
        raise ValueError(f"元メッセージが見つかりません: {trigger_message_id}")
    trigger = dict(trigger)

    if client is None:
        from services.chatwork import ChatworkClient
        client = ChatworkClient()

    msgs = related_messages(trigger)
    files_text = extract_files_text(client, trigger["room_id"], msgs)
    prompt = build_prompt(trigger, msgs, files_text)
    parsed, env = run_json(prompt, model=model(), timeout=300)

    body = (parsed.get("body_md") or "").strip()
    if not body:
        raise ValueError("AIが本文を返しませんでした。")

    period = datetime.date.fromtimestamp(trigger["send_time"]).strftime("%Y-%m")
    evidence = [m["message_id"] for m in msgs]
    files_json = [{"filename": fn, "ok": txt is not None, "error": err}
                  for fn, txt, err in files_text]
    save(trigger_message_id, period, trigger["room_id"], body, parsed.get("summary"),
        evidence, files_json, model(), generated_by)
    _log(trigger_message_id, prompt, env, parsed)
    return get_by_trigger(trigger_message_id)


def save(trigger_message_id, period, room_id, body, summary, evidence, files_json,
        model_name, generated_by):
    with get_conn() as conn:
        conn.execute(
            "INSERT INTO monthly_reports (trigger_message_id, report_period, room_id, body, "
            "summary, evidence, files, model, generated_by) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?) "
            "ON CONFLICT(trigger_message_id) DO UPDATE SET "
            "report_period=excluded.report_period, room_id=excluded.room_id, body=excluded.body, "
            "summary=excluded.summary, evidence=excluded.evidence, files=excluded.files, "
            "model=excluded.model, generated_by=excluded.generated_by, updated_at=datetime('now')",
            (trigger_message_id, period, room_id, body, summary,
             json.dumps(evidence, ensure_ascii=False), json.dumps(files_json, ensure_ascii=False),
             model_name, generated_by))


def _log(trigger_message_id, prompt, env, parsed):
    with get_conn() as conn:
        conn.execute(
            "INSERT INTO ai_analysis_logs (kind, model, prompt, raw_output, parsed, duration_ms) "
            "VALUES ('monthly_report', ?, ?, ?, ?, ?)",
            (model(), prompt[:20000], str(env.get("result"))[:20000],
             json.dumps({"trigger_message_id": trigger_message_id, **parsed}, ensure_ascii=False)[:20000],
             env.get("_elapsed_ms")))


def get_by_trigger(trigger_message_id: str):
    r = query_one("SELECT * FROM monthly_reports WHERE trigger_message_id=?", (trigger_message_id,))
    return dict(r) if r else None


def list_all(limit: int = 24):
    return [dict(r) for r in query(
        "SELECT * FROM monthly_reports ORDER BY report_period DESC, id DESC LIMIT ?", (limit,))]


def delete(trigger_message_id: str):
    with get_conn() as conn:
        conn.execute("DELETE FROM monthly_reports WHERE trigger_message_id=?", (trigger_message_id,))
