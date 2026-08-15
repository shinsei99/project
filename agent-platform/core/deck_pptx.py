"""講演・提案スライドを PowerPoint（16:9）で組む

なぜ作り直したか:
  前の作りは**全ページ同じ型**（左に「・」箇条書き、右に画像）で、
  12枚並べても単調だった。講演では、章の切り替え・数字・引用・図版で
  面の性格を変えないと、聞き手の集中が続かない。

  出力は .pptx にする。**講演はPowerPointで行うため**、
  当日その場で文字を直せることが最優先。PDFでは直せない。

崩れさせないための決まり:
  - 文字は**必ず枠の中に収める**。文字数から級数を決め、収まらなければ落とす
  - 画像は**縦横比を保って**枠に収める（引き伸ばさない）
  - 1枚に置く要素は最大6つまで。詰め込むと読めない
"""
from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

# 用紙。**既定は4:3**。
# 実際に使われている講演資料（Googleドライブの16本）を調べたところ、
# すべて 25.4 x 19.05cm（4:3）だった。会場のプロジェクタに合わせた形なので、
# こちらに合わせる。16:9 が要るときは paper="16:9" を渡す。
PAPERS = {"4:3": (25.4, 19.05), "16:9": (33.867, 19.05)}
W = 25.4
H = 19.05
MARGIN = 1.6

INK = "#12243d"
ACCENT = "#e0533a"
SOFT = "#f2f5f8"
GRAY = "#5b6675"

# 日本語のフォント。**いつもの資料は丸ゴシック**（HG丸ｺﾞｼｯｸM-PRO）だったので、
# Macに入っている丸ゴシックを既定にする。無い環境ではヒラギノ角ゴに落ちる。
FONT = "ヒラギノ丸ゴ ProN W4"
FONT_BOLD = "ヒラギノ丸ゴ ProN W4"


def _cm(value: float):
    from pptx.util import Cm

    return Cm(value)


def _rgb(color: str):
    from pptx.dml.color import RGBColor

    text = str(color).lstrip("#")
    return RGBColor(int(text[0:2], 16), int(text[2:4], 16), int(text[4:6], 16))


def _fit_size(text: str, base: int, per_line: int, min_size: int) -> int:
    """文字数に応じて級数を落とす。**枠からはみ出させない**ための保険。"""
    length = len(str(text or ""))
    if length <= per_line:
        return base
    shrink = base * per_line / float(length)
    return int(max(min_size, shrink))


def _image_box(path, box: Tuple[float, float, float, float]):
    """枠(left, top, width, height)に、縦横比を保って収める寸法を返す。"""
    left, top, width, height = box
    try:
        from PIL import Image

        with Image.open(path) as img:
            iw, ih = img.size
    except Exception:
        return left, top, width, height
    if not iw or not ih:
        return left, top, width, height
    scale = min(width / iw, height / ih)
    new_w, new_h = iw * scale, ih * scale
    return (left + (width - new_w) / 2, top + (height - new_h) / 2, new_w, new_h)


class DeckBuilder:
    """スライドを1枚ずつ足していく。"""

    def __init__(self, accent: str = ACCENT, ink: str = INK, paper: str = "4:3"):
        from pptx import Presentation

        global W, H
        W, H = PAPERS.get(str(paper), PAPERS["4:3"])
        self.accent = accent
        self.ink = ink
        self.paper = paper
        self.prs = Presentation()
        self.prs.slide_width = _cm(W)
        self.prs.slide_height = _cm(H)
        self.count = 0

    # --- 部品 ---------------------------------------------------------------

    def _blank(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _rect(self, slide, box, color, line=None):
        from pptx.enum.shapes import MSO_SHAPE

        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _cm(box[0]), _cm(box[1]),
                                       _cm(box[2]), _cm(box[3]))
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(color)
        if line:
            shape.line.color.rgb = _rgb(line)
            shape.line.width = _cm(0.03)
        else:
            shape.line.fill.background()
        shape.shadow.inherit = False
        return shape

    def _text(self, slide, box, text, size, color=INK, bold=False, align="left",
              line_spacing=1.4, anchor="top"):
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Pt

        tb = slide.shapes.add_textbox(_cm(box[0]), _cm(box[1]), _cm(box[2]), _cm(box[3]))
        frame = tb.text_frame
        frame.word_wrap = True
        frame.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
                                 "bottom": MSO_ANCHOR.BOTTOM}[anchor]
        frame.margin_left = frame.margin_right = 0
        frame.margin_top = frame.margin_bottom = 0

        lines = str(text or "").split("\n")
        for index, line in enumerate(lines):
            para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            para.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                              "right": PP_ALIGN.RIGHT}[align]
            para.line_spacing = line_spacing
            run = para.add_run()
            run.text = line
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.name = FONT
            run.font.color.rgb = _rgb(color)
        return tb

    def _picture(self, slide, path, box):
        if not path or not Path(path).exists():
            return None
        left, top, width, height = _image_box(path, box)
        return slide.shapes.add_picture(str(path), _cm(left), _cm(top),
                                        width=_cm(width), height=_cm(height))

    def _picture_cover(self, slide, path, box):
        """枠を**埋める**ように写真を置く（はみ出しは切る）。

        縦横比を保ったまま枠いっぱいに広げ、余った側を切る。
        紙面のcssでいう object-fit: cover と同じ考え方。
        全面写真の面はこれが無いと、余白ができて締まらない。
        """
        if not path or not Path(path).exists():
            return None
        left, top, width, height = box
        try:
            from PIL import Image

            with Image.open(path) as img:
                iw, ih = img.size
        except Exception:
            return self._picture(slide, path, box)
        if not iw or not ih:
            return None
        box_ratio = width / height
        img_ratio = iw / ih
        pic = slide.shapes.add_picture(str(path), _cm(left), _cm(top),
                                       width=_cm(width), height=_cm(height))
        # はみ出す側だけを切る（python-pptx の crop は 0〜1 の割合）
        if img_ratio > box_ratio:
            trim = (1 - box_ratio / img_ratio) / 2
            pic.crop_left = pic.crop_right = trim
        else:
            trim = (1 - img_ratio / box_ratio) / 2
            pic.crop_top = pic.crop_bottom = trim
        return pic

    def _scrim(self, slide, box, color: str, alpha: int = 55):
        """写真の上に敷く半透明の面。**文字を読ませるために要る**。

        python-pptx に透明度の指定が無いので、XMLに直接 alpha を書く。
        """
        from pptx.oxml.ns import qn

        shape = self._rect(slide, box, color)
        fill = shape.fill.fore_color._xFill  # noqa: SLF001
        srgb = fill.find(qn("a:srgbClr"))
        if srgb is not None:
            element = srgb.makeelement(qn("a:alpha"), {"val": str(int(alpha * 1000))})
            srgb.append(element)
        return shape

    def _page_number(self, slide):
        self._text(slide, (W - 3.2, H - 1.5, 2.0, 0.9), "%02d" % self.count, 11,
                   color="#9aa3 b0".replace(" ", ""), align="right")

    def _chrome(self, slide, tag: str = ""):
        """全ページ共通の骨（上の帯・章名・ページ番号）。"""
        self._rect(slide, (0, 0, W, 0.28), self.ink)
        if tag:
            self._text(slide, (MARGIN, 0.95, 12, 0.8), tag, 11, color="#98a2b0",
                       bold=True)
        self._page_number(slide)

    # --- 面の種類 -----------------------------------------------------------

    def cover(self, title: str, lead: str = "", kicker: str = "", who: str = ""):
        self.count += 1
        slide = self._blank()
        self._rect(slide, (0, 0, W, H), self.ink)
        self._rect(slide, (0, 0, W, 0.42), self.accent)
        top = 4.4
        if kicker:
            self._text(slide, (MARGIN + 0.8, top, W - 6, 1.0), kicker, 14,
                       color="#c9d2dd", bold=True)
            top += 1.5
        size = _fit_size(title, 44, 18, 26)
        self._text(slide, (MARGIN + 0.8, top, W - 6, 4.6), title, size, color="#ffffff",
                   bold=True, line_spacing=1.25)
        top += 4.9
        if lead:
            self._text(slide, (MARGIN + 0.8, top, W - 7, 2.4), lead, 14,
                       color="#dbe3ec", line_spacing=1.6)
        if who:
            self._text(slide, (MARGIN + 0.8, H - 4.0, W - 7, 2.4), who, 12,
                       color="#b8c3d0", line_spacing=1.7)
        return slide

    def section(self, number: str, title: str, lead: str = ""):
        self.count += 1
        slide = self._blank()
        self._rect(slide, (0, 0, W, H), SOFT)
        self._rect(slide, (0, 0, W, 0.28), self.ink)
        self._text(slide, (MARGIN + 0.8, 6.2, W - 6, 0.9), number, 13,
                   color=self.accent, bold=True)
        size = _fit_size(title, 36, 16, 24)
        self._text(slide, (MARGIN + 0.8, 7.3, W - 6, 3.4), title, size, color=self.ink,
                   bold=True, line_spacing=1.3)
        if lead:
            self._text(slide, (MARGIN + 0.8, 11.0, W - 7, 2.0), lead, 13, color=GRAY,
                       line_spacing=1.6)
        self._page_number(slide)
        return slide

    def _heading(self, slide, title: str, sub: str = "", tag: str = ""):
        self._chrome(slide, tag)
        size = _fit_size(title, 28, 22, 18)
        self._text(slide, (MARGIN, 2.0, W - MARGIN * 2, 2.0), title, size,
                   color=self.ink, bold=True, line_spacing=1.25)
        top = 4.2
        if sub:
            self._text(slide, (MARGIN, top, W - MARGIN * 2, 1.4), sub, 12, color=GRAY,
                       line_spacing=1.5)
            top += 1.5
        self._rect(slide, (MARGIN, top, 2.0, 0.16), self.accent)
        return top + 0.9

    def bullets(self, title: str, items: List[str], sub: str = "", tag: str = "",
                image: Optional[str] = None, caption: str = ""):
        """要点。画像があれば右に置く。"""
        self.count += 1
        slide = self._blank()
        top = self._heading(slide, title, sub, tag)

        text_width = (W - MARGIN * 2) * (0.52 if image else 1.0)
        for index, item in enumerate(items[:5]):
            y = top + index * 1.75
            self._rect(slide, (MARGIN, y + 0.34, 0.34, 0.34), self.accent)
            self._text(slide, (MARGIN + 0.85, y, text_width - 1.0, 1.6),
                       str(item), 15, color="#22303f", line_spacing=1.45)

        if image:
            box = (MARGIN + text_width + 0.8, top,
                   W - MARGIN * 2 - text_width - 0.8, H - top - 2.2)
            self._picture(slide, image, box)
            if caption:
                self._text(slide, (box[0], H - 2.0, box[2], 0.8), caption, 10,
                           color="#8b95a4")
        return slide

    def figure(self, title: str, image: str, sub: str = "", caption: str = "",
               tag: str = ""):
        """図版を大きく1点。"""
        self.count += 1
        slide = self._blank()
        top = self._heading(slide, title, sub, tag)
        self._rect(slide, (MARGIN, top, W - MARGIN * 2, H - top - 2.2), SOFT)
        self._picture(slide, image, (MARGIN + 0.5, top + 0.5,
                                     W - MARGIN * 2 - 1.0, H - top - 3.2))
        if caption:
            self._text(slide, (MARGIN, H - 1.9, W - MARGIN * 2, 0.8), caption, 10,
                       color="#8b95a4")
        return slide

    def photo(self, title: str, image: str, items: List[str], caption: str = "",
              tag: str = ""):
        """写真を左に大きく、要点を右に。"""
        self.count += 1
        slide = self._blank()
        top = self._heading(slide, title, "", tag)
        half = (W - MARGIN * 2 - 1.2) / 2
        self._picture(slide, image, (MARGIN, top, half, H - top - 2.4))
        if caption:
            self._text(slide, (MARGIN, H - 2.0, half, 0.8), caption, 10,
                       color="#8b95a4")
        for index, item in enumerate(items[:4]):
            y = top + index * 1.9
            self._rect(slide, (MARGIN + half + 1.2, y + 0.36, 0.32, 0.32), self.accent)
            self._text(slide, (MARGIN + half + 2.0, y, half - 0.8, 1.7), str(item), 15,
                       color="#22303f", line_spacing=1.45)
        return slide

    def numbers(self, title: str, items: List[Dict[str, str]], note: str = "",
                tag: str = ""):
        """数字を大きく並べる。講演で一番残るのは数字。"""
        self.count += 1
        slide = self._blank()
        top = self._heading(slide, title, "", tag)
        items = items[:4]
        gap = 0.7
        width = (W - MARGIN * 2 - gap * (len(items) - 1)) / max(len(items), 1)
        height = 6.4
        for index, item in enumerate(items):
            left = MARGIN + index * (width + gap)
            self._rect(slide, (left, top, width, height), SOFT)
            value = str(item.get("value", ""))
            unit = str(item.get("unit", ""))
            self._text(slide, (left, top + 1.2, width, 2.6),
                       value, _fit_size(value, 40, 4, 22), color=self.accent,
                       bold=True, align="center")
            if unit:
                self._text(slide, (left, top + 3.7, width, 0.9), unit, 13,
                           color=self.accent, bold=True, align="center")
            self._text(slide, (left + 0.5, top + 4.6, width - 1.0, 1.6),
                       str(item.get("label", "")), 11, color=GRAY, align="center",
                       line_spacing=1.4)
        if note:
            self._text(slide, (MARGIN, top + height + 0.9, W - MARGIN * 2, 1.6),
                       note, 13, color="#22303f", line_spacing=1.6)
        return slide

    def icons(self, title: str, items: List[Dict[str, str]], sub: str = "",
              tag: str = ""):
        """記号つきの3〜4点。文字だけの面が続くのを避けるために使う。"""
        import tools

        self.count += 1
        slide = self._blank()
        top = self._heading(slide, title, sub, tag)
        items = items[:4]
        gap = 0.7
        width = (W - MARGIN * 2 - gap * (len(items) - 1)) / max(len(items), 1)
        height = H - top - 2.2
        for index, item in enumerate(items):
            left = MARGIN + index * (width + gap)
            self._rect(slide, (left, top, width, height), SOFT)
            icon = None
            try:
                icon = tools.symbols.png(item.get("icon", ""), self.accent, 256)
            except Exception:
                icon = None
            text_top = top + 1.0
            if icon:
                self._picture(slide, icon, (left + width / 2 - 1.1, top + 0.9, 2.2, 2.2))
                text_top = top + 3.6
            self._text(slide, (left + 0.6, text_top, width - 1.2, 1.6),
                       str(item.get("title", "")), 15, color=self.ink, bold=True,
                       align="center", line_spacing=1.35)
            self._text(slide, (left + 0.6, text_top + 1.8, width - 1.2, height - 4.8),
                       str(item.get("text", "")), 11, color=GRAY, align="center",
                       line_spacing=1.5)
        return slide

    def hero(self, title: str, image: str, sub: str = "", kicker: str = "",
             tag: str = ""):
        """**写真を全面に敷いて、その上に文字**。章の切り替えや表紙に効く。

        いつもの講演資料が写真主体だったのに合わせた面。
        文字が読めるよう、写真の上に濃い色の膜を敷く。
        """
        self.count += 1
        slide = self._blank()
        if image and Path(image).exists():
            self._picture_cover(slide, image, (0, 0, W, H))
            self._scrim(slide, (0, 0, W, H), self.ink, alpha=52)
        else:
            self._rect(slide, (0, 0, W, H), self.ink)
        self._rect(slide, (0, 0, W, 0.42), self.accent)
        top = H * 0.42
        if kicker:
            self._text(slide, (MARGIN + 0.6, top - 1.5, W - 4, 1.0), kicker, 13,
                       color="#e6ecf3", bold=True)
        self._text(slide, (MARGIN + 0.6, top, W - 4, 3.4), title,
                   _fit_size(title, 34, 16, 22), color="#ffffff", bold=True,
                   line_spacing=1.25)
        if sub:
            self._text(slide, (MARGIN + 0.6, top + 3.6, W - 5, 2.2), sub, 13,
                       color="#dfe6ee", line_spacing=1.6)
        self._page_number(slide)
        return slide

    def photos(self, title: str, images: List[str], captions: List[str] = (),
               sub: str = "", tag: str = ""):
        """**写真を敷き詰める面。** いつもの講演資料はこれが主役。

        実際の資料は1枚に1〜6点の写真が入り、文字は見出しだけだった。
        文字で説明するより、現場の写真と新聞記事を並べたほうが伝わる。
        """
        self.count += 1
        slide = self._blank()
        top = self._heading(slide, title, sub, tag)
        images = [x for x in (images or []) if x and Path(x).exists()][:6]
        if not images:
            return slide

        cols = 1 if len(images) == 1 else (2 if len(images) <= 4 else 3)
        rows = (len(images) + cols - 1) // cols
        gap = 0.4
        area_h = H - top - 1.4
        cell_w = (W - MARGIN * 2 - gap * (cols - 1)) / cols
        cell_h = (area_h - gap * (rows - 1)) / rows
        for index, image in enumerate(images):
            row, col = divmod(index, cols)
            box = (MARGIN + col * (cell_w + gap), top + row * (cell_h + gap),
                   cell_w, cell_h - (0.55 if captions else 0))
            self._picture(slide, image, box)
            if captions and index < len(captions) and captions[index]:
                self._text(slide, (box[0], box[1] + box[3] + 0.05, cell_w, 0.5),
                           str(captions[index]), 9, color="#8b95a4", align="center")
        return slide

    def quote(self, text: str, source: str = "", tag: str = ""):
        """大きな一文。講演の山場に置く。"""
        self.count += 1
        slide = self._blank()
        self._chrome(slide, tag)
        self._rect(slide, (MARGIN, 5.4, 0.22, 6.0), self.accent)
        size = _fit_size(text, 30, 26, 18)
        self._text(slide, (MARGIN + 1.2, 5.4, W - MARGIN * 2 - 1.6, 5.2), text, size,
                   color=self.ink, bold=True, line_spacing=1.5, anchor="middle")
        if source:
            self._text(slide, (MARGIN + 1.2, 11.4, W - MARGIN * 2 - 1.6, 2.0),
                       source, 12, color=GRAY, line_spacing=1.6)
        return slide

    def closing(self, title: str, lines: List[str]):
        self.count += 1
        slide = self._blank()
        self._rect(slide, (0, 0, W, H), self.ink)
        self._rect(slide, (0, 0, W, 0.42), self.accent)
        self._text(slide, (MARGIN + 0.8, 4.6, W - 6, 2.4),
                   title, _fit_size(title, 32, 18, 22), color="#ffffff", bold=True)
        body = "\n".join(str(x) for x in lines[:5])
        self._text(slide, (MARGIN + 0.8, 8.2, W - 6, 8.0), body, 14, color="#dbe3ec",
                   line_spacing=1.9)
        return slide

    # --- 出力 ---------------------------------------------------------------

    def save(self, path) -> Path:
        path = Path(path)
        path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(path))
        return path


def build(spec: List[Dict[str, Any]], out_path, accent: str = ACCENT,
          ink: str = INK, paper: str = "4:3") -> Path:
    """面の指定（spec）からスライドを組む。

    spec は [{"type": "cover", ...}, {"type": "bullets", ...}] の配列。
    知らない type は bullets として扱う（1枚も落とさないため）。
    """
    deck = DeckBuilder(accent=accent, ink=ink, paper=paper)
    for item in spec or []:
        if not isinstance(item, dict):
            continue
        kind = str(item.get("type", "bullets"))
        try:
            if kind == "hero":
                deck.hero(item.get("title", ""), item.get("image", ""),
                          item.get("sub", ""), item.get("kicker", ""),
                          item.get("tag", ""))
            elif kind == "cover":
                deck.cover(item.get("title", ""), item.get("lead", ""),
                           item.get("kicker", ""), item.get("who", ""))
            elif kind == "section":
                deck.section(item.get("number", ""), item.get("title", ""),
                             item.get("lead", ""))
            elif kind == "figure":
                deck.figure(item.get("title", ""), item.get("image", ""),
                            item.get("sub", ""), item.get("caption", ""),
                            item.get("tag", ""))
            elif kind == "photos":
                deck.photos(item.get("title", ""), item.get("images", []),
                            item.get("captions", []), item.get("sub", ""),
                            item.get("tag", ""))
            elif kind == "photo":
                deck.photo(item.get("title", ""), item.get("image", ""),
                           item.get("items", []), item.get("caption", ""),
                           item.get("tag", ""))
            elif kind == "numbers":
                deck.numbers(item.get("title", ""), item.get("items", []),
                             item.get("note", ""), item.get("tag", ""))
            elif kind == "icons":
                deck.icons(item.get("title", ""), item.get("items", []),
                           item.get("sub", ""), item.get("tag", ""))
            elif kind == "quote":
                deck.quote(item.get("text", ""), item.get("source", ""),
                           item.get("tag", ""))
            elif kind == "closing":
                deck.closing(item.get("title", ""), item.get("lines", []))
            else:
                deck.bullets(item.get("title", ""), item.get("items", []),
                             item.get("sub", ""), item.get("tag", ""),
                             item.get("image"), item.get("caption", ""))
        except Exception:
            continue
    return deck.save(out_path)
